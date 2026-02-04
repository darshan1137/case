import os
from typing import Optional, List, Any, Dict
from dotenv import load_dotenv

# Load environment variables FIRST before any imports
load_dotenv()

# Set USER_AGENT immediately after loading .env
if not os.getenv("USER_AGENT"):
    os.environ["USER_AGENT"] = "AI-Cloud-RAG-Service/1.0"

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field, create_model
import aiofiles
import json
import tempfile
from functools import lru_cache

# LangChain & Groq
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEndpointEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import AsyncHtmlLoader, TextLoader
from langchain_community.document_transformers import Html2TextTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.documents import Document

app = FastAPI(title="Groq Universal RAG API", version="1.0.0")

# --- CACHING LAYER ---
class VectorStoreCache:
    """In-memory cache for embeddings and vector stores"""
    def __init__(self):
        self.cache = {}
        self.embeddings_instance = None
    
    def get_embeddings(self):
        """Singleton pattern for embeddings instance"""
        if self.embeddings_instance is None:
            self.embeddings_instance = HuggingFaceEndpointEmbeddings(
                model="sentence-transformers/all-MiniLM-L6-v2",
                huggingfacehub_api_token=os.getenv("HUGGINGFACE_API_KEY")
            )
        return self.embeddings_instance
    
    def get_vectorstore(self, key: str):
        """Retrieve cached vector store by key"""
        return self.cache.get(key)
    
    def set_vectorstore(self, key: str, vectorstore):
        """Cache a vector store with expiry consideration"""
        self.cache[key] = vectorstore
    
    def clear_cache(self):
        """Clear all cached vector stores"""
        self.cache.clear()
    
    def cache_size(self):
        """Get number of cached vector stores"""
        return len(self.cache)

# Initialize global cache
vector_cache = VectorStoreCache()

# Enable CORS for all origins
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- CONFIGURATION ---
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
KNOWLEDGE_BASE_CACHE_KEY = "knowledge_base_faiss"
FILE_CACHE_TTL = 3600  # Cache file-based vector stores for 1 hour

# --- DATA MODELS ---

class RAGRequest(BaseModel):
    source_url: Optional[str] = Field(None, description="URL to scrape")
    source_text: Optional[str] = Field(None, description="Raw text content")
    query: str = Field(..., description="User question")
    system_prompt: str = Field(
        "You are a helpful AI. Answer based strictly on the context provided.",
        description="System instructions"
    )
    response_schema: Optional[str] = Field(
        None, 
        description="JSON Schema as string for structured output. If None, returns string."
    )

class RAGResponse(BaseModel):
    result: Any
    used_sources: List[str]

# --- HELPER FUNCTIONS ---

def load_and_process_data(url: str = None, text: str = None) -> List[Document]:
    docs = []
    
    # 1. Scrape URL (if provided)
    if url:
        try:
            loader = AsyncHtmlLoader([url])
            raw_docs = loader.load()
            # Convert HTML to clean Markdown
            html2text = Html2TextTransformer()
            docs.extend(html2text.transform_documents(raw_docs))
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Scraping failed: {str(e)}")

    # 2. Load Raw Text (if provided)
    if text:
        docs.append(Document(page_content=text, metadata={"source": "raw_input"}))

    if not docs:
        raise HTTPException(status_code=400, detail="No source provided.")

    # 3. Split Text
    splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=200)
    return splitter.split_documents(docs)

# --- CORE ENDPOINT ---

@app.post("/generate", response_model=RAGResponse)
async def generate_rag(request: RAGRequest):
    # 1. Check API Key
    if not os.getenv("GROQ_API_KEY"):
        raise HTTPException(status_code=500, detail="GROQ_API_KEY is missing.")

    # 2. Ingest & Split
    splits = load_and_process_data(request.source_url, request.source_text)

    # 3. Embed & Index (Ephemeral / In-Memory with Caching)
    # Reuse embeddings instance to avoid reinitialization
    embeddings = vector_cache.get_embeddings()
    vectorstore = FAISS.from_documents(splits, embeddings)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 5})

    # 4. Retrieve Context
    retrieved_docs = retriever.invoke(request.query)
    context_text = "\n\n".join([d.page_content for d in retrieved_docs])

    # 5. Initialize Groq LLM
    # 'mixtral-8x7b-32768' or 'llama3-70b-8192' are great free options on Groq
    llm = ChatGroq(
        temperature=0,
        model_name="openai/gpt-oss-120b", 
        api_key=os.getenv("GROQ_API_KEY")
    )

    # 6. Generate Response
    
    if request.response_schema:
        # --- STRUCTURED MODE (JSON) ---
        try:
            # Parse schema if it's a JSON string
            if isinstance(request.response_schema, str):
                schema = json.loads(request.response_schema)
            else:
                schema = request.response_schema
            
            # Validate schema has required fields
            if not isinstance(schema, dict) or 'title' not in schema or 'properties' not in schema:
                raise ValueError("Schema must have 'title' and 'properties' fields")
            
            prompt = ChatPromptTemplate.from_messages([
                ("system", "{system_prompt}"),
                ("human", "Context: {context}\n\nQuestion: {query}")
            ])
            
            # Bind the schema to the LLM
            structured_llm = llm.with_structured_output(schema)
            chain = prompt | structured_llm
            
            response = chain.invoke({
                "system_prompt": request.system_prompt,
                "context": context_text,
                "query": request.query
            })
        except Exception as e:
            # Fallback to text mode if structured output fails
            response = f"Error in structured mode: {str(e)}. Returning as text instead."

    else:
        # --- STANDARD TEXT MODE ---
        prompt = ChatPromptTemplate.from_messages([
            ("system", "{system_prompt}"),
            ("human", "Context: {context}\n\nQuestion: {query}")
        ])
        
        chain = prompt | llm
        response_msg = chain.invoke({
            "system_prompt": request.system_prompt,
            "context": context_text,
            "query": request.query
        })
        response = response_msg.content

    # Cleanup (note: embeddings instance is kept for reuse)
    del vectorstore

    return RAGResponse(
        result=response,
        used_sources=[d.metadata.get("source", "unknown") for d in retrieved_docs]
    )

# --- HELPER FUNCTION FOR FILE PROCESSING ---

def extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from various file types"""
    file_ext = os.path.splitext(filename)[1].lower()
    
    # Text-based files
    if file_ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.log']:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
    
    # PDF files
    elif file_ext == '.pdf':
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() for page in reader.pages)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PDF extraction failed: {str(e)}. Install PyPDF2: pip install PyPDF2")
    
    # Word documents
    elif file_ext in ['.docx', '.doc']:
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"DOCX extraction failed: {str(e)}. Install python-docx: pip install python-docx")
    
    # Images (OCR)
    elif file_ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
        try:
            from PIL import Image
            import pytesseract
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            if not text.strip():
                raise ValueError("No text found in image")
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Image OCR failed: {str(e)}. Install: pip install pillow pytesseract")
    
    # PowerPoint
    elif file_ext in ['.pptx', '.ppt']:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PowerPoint extraction failed: {str(e)}. Install: pip install python-pptx")
    
    # Excel
    elif file_ext in ['.xlsx', '.xls']:
        try:
            import pandas as pd
            df = pd.read_excel(file_path)
            return df.to_string()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Excel extraction failed: {str(e)}. Install: pip install pandas openpyxl")
    
    else:
        # Try to read as text anyway
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")

@app.post("/upload-and-query")
async def upload_and_query(
    file: UploadFile = File(..., description="File to analyze (PDF, DOCX, TXT, images, etc.)"),
    query: str = Form(..., description="Your question about the file"),
    system_prompt: str = Form(
        default="You are a helpful AI. Answer based strictly on the context provided.",
        description="Custom system instructions"
    ),
    response_schema: Optional[str] = Form(None, description="JSON Schema as string for structured output"),
    use_cache: Optional[bool] = Form(True, description="Whether to use cached vector store")
):
    """
    Upload ANY file type and query its contents with optional caching.
    
    Supported formats:
    - Documents: PDF, DOCX, TXT, MD
    - Images: PNG, JPG (OCR)
    - Spreadsheets: XLSX, XLS, CSV
    - Presentations: PPTX
    - Code/Data: JSON, XML, HTML, LOG
    
    Caching: For frequently queried files (e.g., knowledge_base.md), set use_cache=true
    to reuse the vector store across multiple requests.
    """
    if not os.getenv("GROQ_API_KEY"):
        raise HTTPException(status_code=500, detail="GROQ_API_KEY is missing.")
    
    # Generate cache key based on filename
    cache_key = f"{file.filename}_faiss"
    
    # Check cache first (if enabled)
    if use_cache:
        cached_vectorstore = vector_cache.get_vectorstore(cache_key)
        if cached_vectorstore:
            print(f"✓ Using cached vector store for {file.filename}")
            retriever = cached_vectorstore.as_retriever(search_kwargs={"k": 5})
            retrieved_docs = retriever.invoke(query)
            context_text = "\n\n".join([d.page_content for d in retrieved_docs])
            
            # Generate response using cached index
            llm = ChatGroq(
                temperature=0,
                model_name="openai/gpt-oss-120b",
                api_key=os.getenv("GROQ_API_KEY")
            )
            
            # Generate Response
            if response_schema:
                try:
                    schema = json.loads(response_schema) if isinstance(response_schema, str) else response_schema
                    if not isinstance(schema, dict) or 'title' not in schema or 'properties' not in schema:
                        raise ValueError("Schema must have 'title' and 'properties' fields")
                    
                    structured_llm = llm.with_structured_output(schema)
                    prompt = ChatPromptTemplate.from_messages([
                        ("system", "{system_prompt}"),
                        ("human", "Context: {context}\n\nQuestion: {query}")
                    ])
                    chain = prompt | structured_llm
                    response = chain.invoke({
                        "system_prompt": system_prompt,
                        "context": context_text,
                        "query": query
                    })
                except Exception as e:
                    prompt = ChatPromptTemplate.from_messages([
                        ("system", "{system_prompt}"),
                        ("human", "Context: {context}\n\nQuestion: {query}")
                    ])
                    chain = prompt | llm
                    response_msg = chain.invoke({
                        "system_prompt": system_prompt,
                        "context": context_text,
                        "query": query
                    })
                    response = f"Note: Structured output failed ({str(e)}). Plain text response:\n{response_msg.content}"
            else:
                prompt = ChatPromptTemplate.from_messages([
                    ("system", "{system_prompt}"),
                    ("human", "Context: {context}\n\nQuestion: {query}")
                ])
                chain = prompt | llm
                response_msg = chain.invoke({
                    "system_prompt": system_prompt,
                    "context": context_text,
                    "query": query
                })
                response = response_msg.content
            
            return RAGResponse(
                result=response,
                used_sources=[f"{file.filename} (cached)"]
            )
    
    # If not cached, process file normally
    print(f"⏳ Processing {file.filename} (not cached)")
    
    # Check file size (max 25MB)
    content = await file.read()
    if len(content) > 25 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large (max 25MB)")
    
    temp_file = None
    try:
        # Save file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name
        
        # Extract text
        text = extract_text_from_file(temp_file_path, file.filename)
        
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from file")
        
        # Create document
        docs = [Document(page_content=text, metadata={"source": file.filename})]
        
        # Split text
        splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=200)
        splits = splitter.split_documents(docs)
        
        # Embed & Index with cached embeddings instance
        embeddings = vector_cache.get_embeddings()
        vectorstore = FAISS.from_documents(splits, embeddings)
        
        # Cache the vector store if enabled
        if use_cache:
            vector_cache.set_vectorstore(cache_key, vectorstore)
            print(f"✓ Cached vector store for {file.filename}. Total cached: {vector_cache.cache_size()}")
        
        retriever = vectorstore.as_retriever(search_kwargs={"k": 5})
        
        # Retrieve Context
        retrieved_docs = retriever.invoke(query)
        context_text = "\n\n".join([d.page_content for d in retrieved_docs])
        
        # Initialize LLM
        llm = ChatGroq(
            temperature=0,
            model_name="openai/gpt-oss-120b",
            api_key=os.getenv("GROQ_API_KEY")
        )
        
        # Generate Response
        if response_schema:
            try:
                schema = json.loads(response_schema) if isinstance(response_schema, str) else response_schema
                if not isinstance(schema, dict) or 'title' not in schema or 'properties' not in schema:
                    raise ValueError("Schema must have 'title' and 'properties' fields")
                
                structured_llm = llm.with_structured_output(schema)
                prompt = ChatPromptTemplate.from_messages([
                    ("system", "{system_prompt}"),
                    ("human", "Context: {context}\n\nQuestion: {query}")
                ])
                chain = prompt | structured_llm
                response = chain.invoke({
                    "system_prompt": system_prompt,
                    "context": context_text,
                    "query": query
                })
            except Exception as e:
                # Fallback to text if structured fails
                prompt = ChatPromptTemplate.from_messages([
                    ("system", "{system_prompt}"),
                    ("human", "Context: {context}\n\nQuestion: {query}")
                ])
                chain = prompt | llm
                response_msg = chain.invoke({
                    "system_prompt": system_prompt,
                    "context": context_text,
                    "query": query
                })
                response = f"Note: Structured output failed ({str(e)}). Plain text response:\n{response_msg.content}"
        else:
            prompt = ChatPromptTemplate.from_messages([
                ("system", "{system_prompt}"),
                ("human", "Context: {context}\n\nQuestion: {query}")
            ])
            chain = prompt | llm
            response_msg = chain.invoke({
                "system_prompt": system_prompt,
                "context": context_text,
                "query": query
            })
            response = response_msg.content
        
        # Cleanup temporary file
        if not use_cache:
            del vectorstore  # Only delete if not caching
        os.remove(temp_file_path)
        
        return RAGResponse(
            result=response,
            used_sources=[file.filename]
        )
        
    except HTTPException:
        raise
    except Exception as e:
        if temp_file and os.path.exists(temp_file.name):
            os.remove(temp_file.name)
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

@app.get("/health")
def health():
    """Health check with cache status"""
    return {
        "status": "active",
        "provider": "groq",
        "cache": {
            "cached_vectorstores": vector_cache.cache_size(),
            "embeddings_initialized": vector_cache.embeddings_instance is not None
        }
    }

@app.post("/cache/clear")
def clear_cache():
    """Clear all cached vector stores (use with caution)"""
    vector_cache.clear_cache()
    return {
        "status": "success",
        "message": "All cached vector stores cleared",
        "cached_count": 0
    }

@app.get("/cache/status")
def cache_status():
    """Get detailed cache status"""
    return {
        "cached_vectorstores": vector_cache.cache_size(),
        "embeddings_model": EMBEDDING_MODEL,
        "embeddings_initialized": vector_cache.embeddings_instance is not None,
        "cache_keys": list(vector_cache.cache.keys())
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(
        "rag:app",
        host="0.0.0.0",
        port=8002,
        reload=True,
        log_level="info"
    )
